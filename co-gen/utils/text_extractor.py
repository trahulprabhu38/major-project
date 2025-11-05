"""
Text extraction utilities for various document formats
"""
import os
import re
from pathlib import Path
from typing import Optional

import pdfplumber
from docx import Document
from pptx import Presentation
import PyPDF2


class TextExtractor:
    """Extract text from various document formats"""

    SUPPORTED_EXTENSIONS = ['.pdf', '.docx', '.pptx', '.txt']

    @staticmethod
    def extract_from_pdf(file_path: str) -> str:
        """Extract text from PDF using pdfplumber"""
        texts = []
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        texts.append(text)
        except Exception as e:
            print(f"pdfplumber failed for {file_path}: {e}")
            # Fallback to PyPDF2
            try:
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        text = page.extract_text()
                        if text:
                            texts.append(text)
            except Exception as e2:
                print(f"PyPDF2 also failed for {file_path}: {e2}")
                raise ValueError(f"Failed to extract text from PDF: {e2}")

        return '\n'.join(texts)

    @staticmethod
    def extract_from_docx(file_path: str) -> str:
        """Extract text from DOCX"""
        try:
            doc = Document(file_path)
            return '\n'.join(paragraph.text for paragraph in doc.paragraphs)
        except Exception as e:
            raise ValueError(f"Failed to extract text from DOCX: {e}")

    @staticmethod
    def extract_from_pptx(file_path: str) -> str:
        """Extract text from PPTX"""
        try:
            prs = Presentation(file_path)
            slides_text = []
            for slide in prs.slides:
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
                slides_text.append('\n'.join(texts))
            return '\n'.join(slides_text)
        except Exception as e:
            raise ValueError(f"Failed to extract text from PPTX: {e}")

    @staticmethod
    def extract_from_txt(file_path: str) -> str:
        """Extract text from TXT"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            raise ValueError(f"Failed to extract text from TXT: {e}")

    @classmethod
    def extract(cls, file_path: str) -> str:
        """Extract text from any supported file format"""
        ext = Path(file_path).suffix.lower()

        if ext not in cls.SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file type: {ext}")

        if ext == '.pdf':
            return cls.extract_from_pdf(file_path)
        elif ext == '.docx':
            return cls.extract_from_docx(file_path)
        elif ext == '.pptx':
            return cls.extract_from_pptx(file_path)
        elif ext == '.txt':
            return cls.extract_from_txt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    @staticmethod
    def extract_syllabus_section(text: str) -> str:
        """
        Extract syllabus section from document text.
        Looks for keywords like SYLLABUS, UNIT, MODULE
        """
        upper_text = text.upper()

        # Keywords to identify start of syllabus
        start_keywords = [
            "SYLLABUS",
            "COURSE CONTENT",
            "UNIT I",
            "UNIT 1",
            "MODULE I",
            "MODULE 1"
        ]

        # Keywords to identify end of syllabus
        end_keywords = [
            "TEXT BOOKS",
            "TEXTBOOKS",
            "REFERENCE BOOKS",
            "REFERENCES",
            "COURSE OUTCOMES",
            "CO1",
            "EXAMINATION",
            "EXAM"
        ]

        # Find start index
        start_idx = -1
        for keyword in start_keywords:
            idx = upper_text.find(keyword)
            if idx != -1:
                start_idx = idx
                break

        if start_idx == -1:
            # No specific keyword found, return full text
            return text.strip()

        # Find end index
        end_idx = len(text)
        for keyword in end_keywords:
            idx = upper_text.find(keyword, start_idx)
            if idx != -1:
                end_idx = min(end_idx, idx)

        return text[start_idx:end_idx].strip()

    @staticmethod
    def chunk_text(text: str, chunk_size: int = 400, overlap: int = 80) -> list[str]:
        """
        Split text into overlapping chunks
        """
        chunks = []
        i = 0
        while i < len(text):
            chunk = text[i:i + chunk_size].strip()
            if chunk and len(chunk) > 20:  # Only keep meaningful chunks
                chunks.append(chunk)
            i += chunk_size - overlap

        return chunks
