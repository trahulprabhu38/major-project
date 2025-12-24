import os
from pathlib import Path
from pdfminer.high_level import extract_text
from pptx import Presentation
from pdf2image import convert_from_path
import pytesseract

RAW_DIR = "data/raw"
OUT_DIR = "data/extracted"

os.makedirs(OUT_DIR, exist_ok=True)

def extract_pdf(path):
    text = extract_text(path)
    if text.strip():
        return text
    
    # fallback: OCR scanned PDFs
    pages = convert_from_path(path, dpi=200)
    full_text = ""
    for p in pages:
        full_text += pytesseract.image_to_string(p) + "\n"
    return full_text


def extract_ppt(path):
    prs = Presentation(path)
    total = []
    for slide in prs.slides:
        temp = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                temp.append(shape.text)
        slide_text = "\n".join(temp)
        total.append(slide_text)
    return "\n\n".join(total)


def extract_all():
    pdf_dir = Path(RAW_DIR) / "pdfs"
    ppt_dir = Path(RAW_DIR) / "syllabus"

    for pdf in pdf_dir.glob("*.pdf"):
        out_path = Path(OUT_DIR) / (pdf.stem + ".txt")
        print("Extracting PDF:", pdf)
        text = extract_pdf(str(pdf))
        out_path.write_text(text, encoding="utf-8")

    for ppt in ppt_dir.glob("*.pptx"):
        out_path = Path(OUT_DIR) / (ppt.stem + ".txt")
        print("Extracting PPT:", ppt)
        text = extract_ppt(str(ppt))
        out_path.write_text(text, encoding="utf-8")

    print("\nDONE — Extracted to data/extracted/")


if __name__ == "__main__":
    extract_all()
