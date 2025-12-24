"""
Document Intelligence Layer
Advanced document processing with semantic chunking and embedding
"""
import re
from typing import List, Dict, Tuple
from sentence_transformers import SentenceTransformer
import numpy as np

class DocumentIntelligence:
    """
    Advanced document processing pipeline:
    - Multi-format extraction (PDF, PPT, DOCX, TXT)
    - Semantic chunking with overlap
    - Embedding generation
    - Metadata extraction
    """
    
    def __init__(self, embedding_model='all-MiniLM-L6-v2'):
        """Initialize with embedding model"""
        self.embedding_model = SentenceTransformer(embedding_model)
        print(f"Document Intelligence Layer initialized with {embedding_model}")
    
    def extract_text(self, file_path: str) -> str:
        """Extract text from various file formats"""
        from PyPDF2 import PdfReader
        from pptx import Presentation
        import docx
        
        text = ""
        file_lower = file_path.lower()
        
        try:
            if file_lower.endswith('.pdf'):
                reader = PdfReader(file_path)
                text = "\n".join([page.extract_text() for page in reader.pages])
            elif file_lower.endswith(('.ppt', '.pptx')):
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            elif file_lower.endswith(('.doc', '.docx')):
                doc = docx.Document(file_path)
                text = "\n".join([para.text for para in doc.paragraphs])
            elif file_lower.endswith('.txt'):
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
        except Exception as e:
            print(f" Error extracting {file_path}: {e}")
        
        return self.clean_text(text)
    
    def clean_text(self, text: str) -> str:
        """Advanced text cleaning"""
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        # Remove special control characters
        text = re.sub(r'[\x00-\x08\x0b-\x0c\x0e-\x1f]', '', text)
        # Remove page numbers and headers
        text = re.sub(r'Page \d+', '', text, flags=re.IGNORECASE)
        # Normalize unicode
        text = text.encode('ascii', 'ignore').decode('ascii')
        return text.strip()
    
    def semantic_chunk(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[Dict]:
        """
        Intelligent semantic chunking with:
        - Sentence boundary detection
        - Topic coherence preservation
        - Overlap for context continuity
        """
        if len(text) <= chunk_size:
            return [{
                'text': text,
                'start': 0,
                'end': len(text),
                'chunk_id': 0
            }]
        
        chunks = []
        sentences = re.split(r'(?<=[.!?])\s+', text)
        current_chunk = []
        current_length = 0
        chunk_id = 0
        
        current_start = 0
        
        for sentence in sentences:
            sentence_length = len(sentence)
            
            if current_length + sentence_length > chunk_size and current_chunk:
                # Save current chunk
                chunk_text = ' '.join(current_chunk)
                chunks.append({
                    'text': chunk_text,
                    'start': current_start,
                    'end': current_start + len(chunk_text),
                    'chunk_id': chunk_id
                })
                chunk_id += 1
                
                # Start new chunk with overlap (last N sentences)
                overlap_sentences = current_chunk[-3:] if len(current_chunk) >= 3 else current_chunk
                overlap_text = ' '.join(overlap_sentences)
                current_start = current_start + len(chunk_text) - len(overlap_text)
                current_chunk = overlap_sentences + [sentence]
                current_length = sum(len(s) for s in current_chunk)
            else:
                current_chunk.append(sentence)
                current_length += sentence_length
        
        # Add final chunk
        if current_chunk:
            chunk_text = ' '.join(current_chunk)
            chunks.append({
                'text': chunk_text,
                'start': current_start,
                'end': current_start + len(chunk_text),
                'chunk_id': chunk_id
            })
        
        return chunks
    
    def generate_embeddings(self, chunks: List[Dict]) -> List[Dict]:
        """Generate embeddings for all chunks"""
        texts = [chunk['text'] for chunk in chunks]
        embeddings = self.embedding_model.encode(texts, show_progress_bar=True)
        
        for i, chunk in enumerate(chunks):
            chunk['embedding'] = embeddings[i].tolist()
            chunk['embedding_dim'] = len(embeddings[i])
        
        return chunks
    
    def extract_metadata(self, text: str, file_path: str) -> Dict:
        """Extract metadata: modules, topics, keywords, structure"""
        metadata = {
            'file_path': file_path,
            'file_name': file_path.split('/')[-1],
            'total_length': len(text),
            'modules': [],
            'topics': [],
            'keywords': []
        }
        
        # Extract modules/units
        module_pattern = r'(?:module|unit|chapter)\s*[0-9]+[:\-]?\s*([^\n]+)'
        modules = re.findall(module_pattern, text, re.IGNORECASE)
        metadata['modules'] = [m.strip() for m in modules[:10]]
        
        # Extract topics (headings)
        lines = text.split('\n')
        for line in lines[:100]:
            line = line.strip()
            if (len(line) > 10 and len(line) < 150 and 
                line[0].isupper() and 
                not line.endswith('.') and
                line.count(' ') < 15):
                metadata['topics'].append(line)
                if len(metadata['topics']) >= 20:
                    break
        
        # Extract keywords (technical terms)
        keywords = ['SQL', 'normalization', 'transaction', 'database', 'schema', 
                   'ER model', 'relational algebra', 'indexing', 'constraints']
        found_keywords = [kw for kw in keywords if kw.lower() in text.lower()]
        metadata['keywords'] = found_keywords
        
        return metadata
    
    def process_document(self, file_path: str) -> Dict:
        """
        Complete document processing pipeline:
        1. Extract text
        2. Clean text
        3. Semantic chunking
        4. Generate embeddings
        5. Extract metadata
        """
        print(f"Processing: {file_path}")
        
        # Extract and clean
        text = self.extract_text(file_path)
        if not text:
            return None
        
        # Chunk
        chunks = self.semantic_chunk(text, chunk_size=1000, overlap=200)
        print(f"    Created {len(chunks)} semantic chunks")
        
        # Embed
        chunks = self.generate_embeddings(chunks)
        print(f"   Generated {len(chunks)} embeddings")
        
        # Metadata
        metadata = self.extract_metadata(text, file_path)
        print(f"   Extracted {len(metadata['modules'])} modules, {len(metadata['topics'])} topics")
        
        return {
            'raw_text': text,
            'chunks': chunks,
            'metadata': metadata,
            'total_chunks': len(chunks)
        }

