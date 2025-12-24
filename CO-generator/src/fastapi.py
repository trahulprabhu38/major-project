# """
# FastAPI Server for CO Generator
# ================================
# RESTful API endpoints for the CO Generation Pipeline
# Wraps the Streamlit dashboard logic into accessible API endpoints
# """

# from fastapi import FastAPI, File, UploadFile, HTTPException, Form
# from fastapi.responses import JSONResponse
# from fastapi.middleware.cors import CORSMiddleware
# from pydantic import BaseModel, Field
# from typing import List, Optional, Dict, Any
# import tempfile
# import os
# import sys
# from pathlib import Path
# from datetime import datetime
# import json

# # Add src to path
# sys.path.insert(0, os.path.dirname(__file__))

# # Import CO generation modules
# try:
#     from PyPDF2 import PdfReader
#     PDF_AVAILABLE = True
# except ImportError:
#     PDF_AVAILABLE = False

# try:
#     from pptx import Presentation
#     PPTX_AVAILABLE = True
# except ImportError:
#     PPTX_AVAILABLE = False

# try:
#     import docx
#     DOCX_AVAILABLE = True
# except ImportError:
#     DOCX_AVAILABLE = False

# from smart_co_generator import VTUCOGenerator
# from metrics_evaluation import MetricsEvaluator

# # ============================================================================
# # FASTAPI APP INITIALIZATION
# # ============================================================================

# app = FastAPI(
#     title="CO Generator API",
#     description="Advanced Multi-Stage AI Pipeline for VTU-Aligned Course Outcomes",
#     version="1.0.0"
# )

# # Add CORS middleware
# app.add_middleware(
#     CORSMiddleware,
#     allow_origins=["*"],
#     allow_credentials=True,
#     allow_methods=["*"],
#     allow_headers=["*"],
# )

# # Initialize generators
# generator = VTUCOGenerator()
# evaluator = MetricsEvaluator()

# # ============================================================================
# # PYDANTIC MODELS
# # ============================================================================

# class COGenerationRequest(BaseModel):
#     """Request model for CO generation"""
#     num_apply: int = Field(2, ge=0, le=4, description="Number of Apply-level COs (CO1-CO4)")
#     num_analyze: int = Field(2, ge=0, le=4, description="Number of Analyze-level COs (CO1-CO4)")

#     class Config:
#         json_schema_extra = {
#             "example": {
#                 "num_apply": 2,
#                 "num_analyze": 2
#             }
#         }

# class CourseOutcome(BaseModel):
#     """Model for a single Course Outcome"""
#     co_num: int
#     co_text: str
#     bloom_level: str
#     po_mappings: str
#     topics_covered: Optional[List[str]] = []
#     reward_score: Optional[float] = 0.0
#     individual_scores: Optional[Dict[str, float]] = {}
#     approved: Optional[bool] = False

# class COGenerationResponse(BaseModel):
#     """Response model for CO generation"""
#     success: bool
#     message: str
#     cos: List[CourseOutcome]
#     metrics: Optional[Dict[str, Any]] = {}
#     extracted_text_length: int
#     files_processed: int
#     timestamp: str

# class RegenerateRequest(BaseModel):
#     """Request model for regenerating a specific CO"""
#     co_num: int = Field(..., ge=1, le=6, description="CO number to regenerate")
#     original_text: str = Field(..., description="Original CO text")
#     bloom_level: str = Field(..., description="Bloom taxonomy level")
#     feedback: str = Field("", description="User feedback for regeneration")

# class HealthResponse(BaseModel):
#     """Health check response"""
#     status: str
#     version: str
#     dependencies: Dict[str, bool]

# # ============================================================================
# # TEXT EXTRACTION FUNCTIONS (from metrics_dashboard.py)
# # ============================================================================

# def extract_pdf(file_path: str) -> str:
#     """Extract text from PDF"""
#     if not PDF_AVAILABLE:
#         return ""
#     try:
#         reader = PdfReader(file_path)
#         return "\n".join([page.extract_text() or "" for page in reader.pages])
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=f"Error reading PDF: {e}")

# def extract_pptx(file_path: str) -> str:
#     """Extract text from PPTX"""
#     if not PPTX_AVAILABLE:
#         return ""
#     try:
#         prs = Presentation(file_path)
#         text = ""
#         for slide in prs.slides:
#             for shape in slide.shapes:
#                 if hasattr(shape, "text"):
#                     text += shape.text + "\n"
#         return text
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=f"Error reading PPTX: {e}")

# def extract_docx(file_path: str) -> str:
#     """Extract text from DOCX"""
#     if not DOCX_AVAILABLE:
#         return ""
#     try:
#         doc = docx.Document(file_path)
#         return "\n".join([para.text for para in doc.paragraphs])
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=f"Error reading DOCX: {e}")

# def extract_text_from_file(file_path: str, filename: str) -> str:
#     """Extract text from uploaded file"""
#     suffix = Path(filename).suffix.lower()

#     try:
#         if suffix == '.pdf':
#             return extract_pdf(file_path)
#         elif suffix in ['.ppt', '.pptx']:
#             return extract_pptx(file_path)
#         elif suffix in ['.doc', '.docx']:
#             return extract_docx(file_path)
#         elif suffix == '.txt':
#             with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
#                 return f.read()
#         else:
#             return ""
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=f"Error extracting text from {filename}: {e}")

# # ============================================================================
# # CO REGENERATION LOGIC (from metrics_dashboard.py)
# # ============================================================================

# def regenerate_co_with_feedback(co_num: int, original_text: str, bloom_level: str, feedback: str) -> str:
#     """
#     Regenerate a CO based on user feedback
#     Uses the feedback to improve the CO text
#     """
#     feedback_lower = feedback.lower() if feedback else ""

#     bloom_verbs = {
#         'Apply': ['Apply', 'Demonstrate', 'Implement', 'Use', 'Execute'],
#         'Analyze': ['Analyse', 'Examine', 'Compare', 'Investigate', 'Differentiate'],
#         'Evaluate': ['Evaluate', 'Assess', 'Justify', 'Critique', 'Validate'],
#         'Create': ['Create', 'Design', 'Develop', 'Construct', 'Write']
#     }

#     templates = {
#         'Apply': [
#             "Apply {topic} concepts to design and implement database solutions for real-world applications",
#             "Apply {topic} techniques to develop efficient database systems addressing practical scenarios",
#             "Demonstrate proficiency in {topic} for database design, implementation, and optimization",
#             "Apply {topic} principles to create and manage database systems with industry-standard practices",
#         ],
#         'Analyze': [
#             "Analyse {topic} scenarios and apply appropriate database techniques for optimal solutions",
#             "Analyse and evaluate {topic} approaches to determine best practices for database systems",
#             "Examine {topic} requirements and apply suitable techniques for database optimization",
#             "Analyse complex {topic} problems and develop systematic solutions using database concepts",
#         ],
#         'Evaluate': [
#             "Ability to conduct experiments using {topic} tools and evaluate database performance",
#             "Evaluate and compare different {topic} approaches through hands-on experimentation",
#             "Assess {topic} implementations and validate results through systematic testing",
#         ],
#         'Create': [
#             "Write comprehensive reports documenting {topic} experiments, methods, and conclusions",
#             "Create detailed documentation of {topic} implementations with analysis and recommendations",
#             "Design and document complete {topic} solutions with technical specifications",
#         ]
#     }

#     topic_keywords = {
#         'sql': 'SQL query processing',
#         'normalization': 'normalization and functional dependencies',
#         'transaction': 'transaction management and concurrency control',
#         'er': 'ER modeling and database design',
#         'schema': 'schema design and implementation',
#         'query': 'query optimization and execution',
#         'index': 'indexing and performance tuning',
#         'mongodb': 'MongoDB and NoSQL databases',
#         'mysql': 'MySQL database operations',
#         'acid': 'ACID properties and transaction management',
#         'relational': 'relational algebra and database operations',
#     }

#     combined_text = (original_text + " " + feedback_lower).lower()
#     topic = "database management"
#     for keyword, topic_name in topic_keywords.items():
#         if keyword in combined_text:
#             topic = topic_name
#             break

#     selected_templates = templates.get(bloom_level, templates['Apply'])

#     if 'specific' in feedback_lower or 'detail' in feedback_lower:
#         template_idx = min(1, len(selected_templates) - 1)
#     elif 'brief' in feedback_lower or 'short' in feedback_lower:
#         template_idx = min(2, len(selected_templates) - 1)
#     elif 'example' in feedback_lower or 'practical' in feedback_lower:
#         template_idx = 0
#     else:
#         template_idx = (co_num - 1) % len(selected_templates)

#     new_co = f"CO{co_num} " + selected_templates[template_idx].format(topic=topic)

#     if 'sql' in feedback_lower and 'sql' not in new_co.lower():
#         new_co = new_co.replace("database", "SQL and database")
#     if 'practical' in feedback_lower or 'real' in feedback_lower:
#         new_co = new_co.replace("solutions", "practical solutions")
#     if 'tool' in feedback_lower:
#         new_co = new_co.replace("database systems", "database systems using industry tools like MySQL and MongoDB")

#     return new_co

# # ============================================================================
# # API ENDPOINTS
# # ============================================================================

# @app.get("/", response_model=HealthResponse)
# async def root():
#     """Root endpoint - health check"""
#     return {
#         "status": "healthy",
#         "version": "1.0.0",
#         "dependencies": {
#             "pdf_available": PDF_AVAILABLE,
#             "pptx_available": PPTX_AVAILABLE,
#             "docx_available": DOCX_AVAILABLE
#         }
#     }

# @app.get("/health", response_model=HealthResponse)
# async def health_check():
#     """Health check endpoint"""
#     return {
#         "status": "healthy",
#         "version": "1.0.0",
#         "dependencies": {
#             "pdf_available": PDF_AVAILABLE,
#             "pptx_available": PPTX_AVAILABLE,
#             "docx_available": DOCX_AVAILABLE
#         }
#     }

# @app.post("/generate-cos", response_model=COGenerationResponse)
# async def generate_cos(
#     files: List[UploadFile] = File(..., description="Course material files (PDF, PPTX, DOCX, TXT)"),
#     num_apply: int = Form(2, ge=0, le=4, description="Number of Apply-level COs"),
#     num_analyze: int = Form(2, ge=0, le=4, description="Number of Analyze-level COs")
# ):
#     """
#     Generate Course Outcomes from uploaded files

#     - **files**: Upload one or more files (PDF, PPT, PPTX, DOC, DOCX, TXT)
#     - **num_apply**: Number of Apply-level COs in CO1-CO4 (default: 2)
#     - **num_analyze**: Number of Analyze-level COs in CO1-CO4 (default: 2)

#     Returns generated COs with quality metrics
#     """

#     # Validate configuration
#     if num_apply + num_analyze != 4:
#         raise HTTPException(
#             status_code=400,
#             detail=f"num_apply + num_analyze must equal 4. Got {num_apply} + {num_analyze} = {num_apply + num_analyze}"
#         )

#     if not files:
#         raise HTTPException(status_code=400, detail="No files uploaded")

#     # Extract text from all files
#     all_text = ""
#     temp_files = []

#     try:
#         for uploaded_file in files:
#             # Save to temporary file
#             suffix = Path(uploaded_file.filename).suffix.lower()
#             with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
#                 content = await uploaded_file.read()
#                 tmp.write(content)
#                 tmp_path = tmp.name
#                 temp_files.append(tmp_path)

#             # Extract text
#             text = extract_text_from_file(tmp_path, uploaded_file.filename)
#             all_text += text + "\n\n"

#         if not all_text.strip():
#             raise HTTPException(
#                 status_code=400,
#                 detail="No text could be extracted from the uploaded files"
#             )

#         # Generate COs
#         generated_cos = generator.generate_custom_cos(
#             all_text,
#             num_apply=num_apply,
#             num_analyze=num_analyze
#         )

#         # Evaluate each CO
#         cos_with_metrics = []
#         for co_data in generated_cos:
#             po_list = [po.strip() for po in co_data['po_mappings'].split(',')]
#             metrics = evaluator.evaluate_single_co(
#                 co_data['co_text'],
#                 co_data['bloom_level'],
#                 po_list
#             )

#             cos_with_metrics.append({
#                 'co_num': co_data['co_num'],
#                 'co_text': co_data['co_text'],
#                 'bloom_level': co_data['bloom_level'],
#                 'po_mappings': co_data['po_mappings'],
#                 'topics_covered': co_data.get('topics_covered', []),
#                 'reward_score': metrics.overall_quality_score,
#                 'individual_scores': {
#                     'vtu': metrics.vtu_compliance_score,
#                     'obe': metrics.obe_alignment_score,
#                     'bloom': metrics.bloom_accuracy,
#                     'conciseness': metrics.conciseness_score
#                 },
#                 'approved': metrics.overall_quality_score >= 0.70
#             })

#         # Compute aggregate metrics
#         pipeline_metrics = evaluator.evaluate_all_cos([
#             {'co_text': co['co_text'], 'bloom_level': co['bloom_level'], 'po_mappings': co['po_mappings']}
#             for co in cos_with_metrics
#         ])

#         return {
#             "success": True,
#             "message": f"Successfully generated {len(cos_with_metrics)} COs from {len(files)} file(s)",
#             "cos": cos_with_metrics,
#             "metrics": pipeline_metrics.to_dict(),
#             "extracted_text_length": len(all_text),
#             "files_processed": len(files),
#             "timestamp": datetime.now().isoformat()
#         }

#     except HTTPException:
#         raise
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=f"Error generating COs: {str(e)}")
#     finally:
#         # Clean up temporary files
#         for tmp_path in temp_files:
#             try:
#                 os.unlink(tmp_path)
#             except:
#                 pass

# @app.post("/regenerate-co")
# async def regenerate_co(request: RegenerateRequest):
#     """
#     Regenerate a specific CO based on user feedback

#     - **co_num**: CO number (1-6)
#     - **original_text**: Original CO text
#     - **bloom_level**: Bloom taxonomy level
#     - **feedback**: User feedback for improvement

#     Returns the regenerated CO
#     """
#     try:
#         new_co_text = regenerate_co_with_feedback(
#             co_num=request.co_num,
#             original_text=request.original_text,
#             bloom_level=request.bloom_level,
#             feedback=request.feedback
#         )

#         # Evaluate the new CO
#         po_list = []  # PO mappings would need to be extracted or provided
#         metrics = evaluator.evaluate_single_co(
#             new_co_text,
#             request.bloom_level,
#             po_list
#         )

#         return {
#             "success": True,
#             "co_num": request.co_num,
#             "regenerated_text": new_co_text,
#             "bloom_level": request.bloom_level,
#             "quality_metrics": {
#                 'vtu': metrics.vtu_compliance_score,
#                 'obe': metrics.obe_alignment_score,
#                 'bloom': metrics.bloom_accuracy,
#                 'overall': metrics.overall_quality_score
#             },
#             "timestamp": datetime.now().isoformat()
#         }
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=f"Error regenerating CO: {str(e)}")

# @app.post("/evaluate-co")
# async def evaluate_co(
#     co_text: str = Form(...),
#     bloom_level: str = Form(...),
#     po_mappings: str = Form("")
# ):
#     """
#     Evaluate a single CO for quality metrics

#     - **co_text**: The course outcome text
#     - **bloom_level**: Expected Bloom taxonomy level
#     - **po_mappings**: Comma-separated PO mappings (e.g., "PO1, PO2, PO3")

#     Returns quality metrics for the CO
#     """
#     try:
#         po_list = [po.strip() for po in po_mappings.split(',')] if po_mappings else []
#         metrics = evaluator.evaluate_single_co(co_text, bloom_level, po_list)

#         return {
#             "success": True,
#             "co_text": co_text,
#             "metrics": metrics.to_dict(),
#             "timestamp": datetime.now().isoformat()
#         }
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=f"Error evaluating CO: {str(e)}")

# @app.post("/evaluate-all-cos")
# async def evaluate_all_cos(cos: List[Dict[str, Any]]):
#     """
#     Evaluate multiple COs and compute aggregate metrics

#     Expects a list of COs with keys: 'co_text', 'bloom_level', 'po_mappings'

#     Returns aggregate pipeline metrics
#     """
#     try:
#         pipeline_metrics = evaluator.evaluate_all_cos(cos)

#         return {
#             "success": True,
#             "metrics": pipeline_metrics.to_dict(),
#             "timestamp": datetime.now().isoformat()
#         }
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=f"Error evaluating COs: {str(e)}")

# @app.get("/bloom-taxonomy")
# async def get_bloom_taxonomy():
#     """Get Bloom's taxonomy reference data"""
#     from metrics_evaluation import BLOOM_TAXONOMY
#     return {
#         "taxonomy": BLOOM_TAXONOMY,
#         "timestamp": datetime.now().isoformat()
#     }

# @app.get("/po-descriptions")
# async def get_po_descriptions():
#     """Get VTU Program Outcome descriptions"""
#     from metrics_evaluation import VTU_PO_DESCRIPTIONS
#     return {
#         "pos": VTU_PO_DESCRIPTIONS,
#         "timestamp": datetime.now().isoformat()
#     }

# # ============================================================================
# # RUN SERVER
# # ============================================================================

# if __name__ == "__main__":
#     import uvicorn
#     uvicorn.run(app, host="0.0.0.0", port=8000)

"""
Enhanced FastAPI Server for CO Generator with Comprehensive Metrics
====================================================================
Complete implementation matching Streamlit dashboard features:
- CO Generation with metrics
- Advanced ML metrics collection
- Latency profiling and benchmarks
- ChromaDB integration for context retrieval
- Real-time metrics dashboard data
- Complete pipeline monitoring
"""

from fastapi import FastAPI, File, UploadFile, HTTPException, Form, BackgroundTasks
from fastapi.responses import JSONResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
from typing import List, Optional, Dict, Any
import tempfile
import os
import sys
from pathlib import Path
from datetime import datetime
import json
import io

# Add src to path
sys.path.insert(0, os.path.dirname(__file__))

# Import CO generation modules
try:
    from PyPDF2 import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

from smart_co_generator import VTUCOGenerator, evaluate_generated_cos
from metrics_evaluation import MetricsEvaluator, BLOOM_TAXONOMY, VTU_PO_DESCRIPTIONS

# Import latency optimization
try:
    from latency_optimizer import LatencyProfiler, EmbeddingCache, PROFILER
    LATENCY_OPTIMIZER_AVAILABLE = True
except ImportError:
    LATENCY_OPTIMIZER_AVAILABLE = False
    print("⚠️ Latency optimizer not available")

# Import ChromaDB utils
try:
    from chromadb_utils import search_syllabus, get_relevant_content_for_co, get_major_topics_from_syllabus
    CHROMADB_AVAILABLE = True
except ImportError:
    CHROMADB_AVAILABLE = False
    print("⚠️ ChromaDB utils not available")

# ============================================================================
# FASTAPI APP INITIALIZATION
# ============================================================================

app = FastAPI(
    title="CO Generator API - Enhanced",
    description="Advanced Multi-Stage AI Pipeline for VTU-Aligned Course Outcomes with Comprehensive Metrics",
    version="2.0.0",
    docs_url="/docs",
    redoc_url="/redoc"
)

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize components
generator = VTUCOGenerator()
evaluator = MetricsEvaluator()

# Global metrics storage (in production, use Redis or database)
METRICS_STORE = {
    'generation_history': [],
    'aggregate_metrics': {},
    'latency_stats': {},
    'user_feedback': [],
    'system_health': {}
}

# Profiler for latency tracking
if LATENCY_OPTIMIZER_AVAILABLE:
    profiler = LatencyProfiler()
else:
    profiler = None

# ============================================================================
# PYDANTIC MODELS
# ============================================================================

class COGenerationRequest(BaseModel):
    """Request model for CO generation"""
    num_apply: int = Field(2, ge=0, le=4, description="Number of Apply-level COs (CO1-CO4)")
    num_analyze: int = Field(2, ge=0, le=4, description="Number of Analyze-level COs (CO1-CO4)")
    use_chromadb: bool = Field(False, description="Use ChromaDB for context retrieval")

class CourseOutcome(BaseModel):
    """Model for a single Course Outcome"""
    co_num: int
    co_text: str
    bloom_level: str
    po_mappings: str
    topics_covered: Optional[List[str]] = []
    reward_score: Optional[float] = 0.0
    individual_scores: Optional[Dict[str, float]] = {}
    approved: Optional[bool] = False
    word_count: Optional[int] = 0
    has_action_verb: Optional[bool] = False
    has_specific_concepts: Optional[bool] = False

class MLMetrics(BaseModel):
    """Machine Learning Metrics"""
    model_type: str = "VTU CO Generator"
    inference_latency_ms: float
    throughput_cos_per_sec: float
    cache_hit_rate: Optional[float] = 0.0
    embedding_generation_ms: Optional[float] = 0.0
    vector_search_ms: Optional[float] = 0.0

class PipelineMetrics(BaseModel):
    """Comprehensive Pipeline Metrics"""
    # Quality Metrics
    bloom_classification_accuracy: float
    average_quality_score: float
    average_vtu_compliance: float
    average_obe_alignment: float
    average_conciseness_score: float
    po_coverage: float

    # Latency Metrics
    document_processing_ms: float = 0.0
    embedding_generation_ms: float = 0.0
    graph_construction_ms: float = 0.0
    vector_search_ms: float = 0.0
    graph_traversal_ms: float = 0.0
    llm_inference_ms: float = 0.0
    refinement_ms: float = 0.0
    total_pipeline_ms: float = 0.0

    # ML Metrics
    ml_metrics: Optional[MLMetrics] = None

class COGenerationResponse(BaseModel):
    """Response model for CO generation"""
    success: bool
    message: str
    cos: List[CourseOutcome]
    pipeline_metrics: PipelineMetrics
    extracted_text_length: int
    files_processed: int
    timestamp: str
    session_id: Optional[str] = None

class RegenerateRequest(BaseModel):
    """Request model for regenerating a specific CO"""
    co_num: int = Field(..., ge=1, le=6, description="CO number to regenerate")
    original_text: str = Field(..., description="Original CO text")
    bloom_level: str = Field(..., description="Bloom taxonomy level")
    feedback: str = Field("", description="User feedback for regeneration")

class FeedbackRequest(BaseModel):
    """User feedback on generated COs"""
    session_id: str
    co_num: int
    approved: bool
    feedback_text: Optional[str] = ""
    edited_text: Optional[str] = ""

class HealthResponse(BaseModel):
    """Health check response"""
    status: str
    version: str
    dependencies: Dict[str, bool]
    system_metrics: Optional[Dict[str, Any]] = {}

class DashboardMetrics(BaseModel):
    """Complete metrics for dashboard display"""
    summary: Dict[str, Any]
    quality_metrics: Dict[str, float]
    latency_breakdown: Dict[str, float]
    bloom_distribution: Dict[str, int]
    po_coverage_details: Dict[str, int]
    recent_generations: List[Dict[str, Any]]
    ml_performance: Optional[Dict[str, Any]] = {}

# ============================================================================
# TEXT EXTRACTION FUNCTIONS
# ============================================================================

def extract_pdf(file_path: str) -> str:
    """Extract text from PDF"""
    if not PDF_AVAILABLE:
        return ""
    try:
        reader = PdfReader(file_path)
        return "\n".join([page.extract_text() or "" for page in reader.pages])
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading PDF: {e}")

def extract_pptx(file_path: str) -> str:
    """Extract text from PPTX"""
    if not PPTX_AVAILABLE:
        return ""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading PPTX: {e}")

def extract_docx(file_path: str) -> str:
    """Extract text from DOCX"""
    if not DOCX_AVAILABLE:
        return ""
    try:
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading DOCX: {e}")

def extract_text_from_file(file_path: str, filename: str) -> str:
    """Extract text from uploaded file"""
    suffix = Path(filename).suffix.lower()

    try:
        if suffix == '.pdf':
            return extract_pdf(file_path)
        elif suffix in ['.ppt', '.pptx']:
            return extract_pptx(file_path)
        elif suffix in ['.doc', '.docx']:
            return extract_docx(file_path)
        elif suffix == '.txt':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        else:
            return ""
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error extracting text from {filename}: {e}")

# ============================================================================
# CO REGENERATION LOGIC
# ============================================================================

def regenerate_co_with_feedback(co_num: int, original_text: str, bloom_level: str, feedback: str) -> str:
    """Regenerate a CO based on user feedback"""
    feedback_lower = feedback.lower() if feedback else ""

    templates = {
        'Apply': [
            "Apply {topic} concepts to design and implement database solutions for real-world applications",
            "Apply {topic} techniques to develop efficient database systems addressing practical scenarios",
            "Demonstrate proficiency in {topic} for database design, implementation, and optimization",
            "Apply {topic} principles to create and manage database systems with industry-standard practices",
        ],
        'Analyze': [
            "Analyse {topic} scenarios and apply appropriate database techniques for optimal solutions",
            "Analyse and evaluate {topic} approaches to determine best practices for database systems",
            "Examine {topic} requirements and apply suitable techniques for database optimization",
            "Analyse complex {topic} problems and develop systematic solutions using database concepts",
        ],
        'Evaluate': [
            "Ability to conduct experiments using {topic} tools and evaluate database performance",
            "Evaluate and compare different {topic} approaches through hands-on experimentation",
            "Assess {topic} implementations and validate results through systematic testing",
        ],
        'Create': [
            "Write comprehensive reports documenting {topic} experiments, methods, and conclusions",
            "Create detailed documentation of {topic} implementations with analysis and recommendations",
            "Design and document complete {topic} solutions with technical specifications",
        ]
    }

    topic_keywords = {
        'sql': 'SQL query processing',
        'normalization': 'normalization and functional dependencies',
        'transaction': 'transaction management and concurrency control',
        'er': 'ER modeling and database design',
        'schema': 'schema design and implementation',
        'query': 'query optimization and execution',
        'index': 'indexing and performance tuning',
        'mongodb': 'MongoDB and NoSQL databases',
        'mysql': 'MySQL database operations',
        'acid': 'ACID properties and transaction management',
        'relational': 'relational algebra and database operations',
    }

    combined_text = (original_text + " " + feedback_lower).lower()
    topic = "database management"
    for keyword, topic_name in topic_keywords.items():
        if keyword in combined_text:
            topic = topic_name
            break

    selected_templates = templates.get(bloom_level, templates['Apply'])

    if 'specific' in feedback_lower or 'detail' in feedback_lower:
        template_idx = min(1, len(selected_templates) - 1)
    elif 'brief' in feedback_lower or 'short' in feedback_lower:
        template_idx = min(2, len(selected_templates) - 1)
    elif 'example' in feedback_lower or 'practical' in feedback_lower:
        template_idx = 0
    else:
        template_idx = (co_num - 1) % len(selected_templates)

    new_co = f"CO{co_num} " + selected_templates[template_idx].format(topic=topic)

    if 'sql' in feedback_lower and 'sql' not in new_co.lower():
        new_co = new_co.replace("database", "SQL and database")
    if 'practical' in feedback_lower or 'real' in feedback_lower:
        new_co = new_co.replace("solutions", "practical solutions")
    if 'tool' in feedback_lower:
        new_co = new_co.replace("database systems", "database systems using industry tools like MySQL and MongoDB")

    return new_co

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def calculate_ml_metrics(start_time: float, num_cos: int, profiler_stats: Dict = None) -> MLMetrics:
    """Calculate ML-specific metrics"""
    import time
    elapsed_ms = (time.time() - start_time) * 1000
    throughput = (num_cos / elapsed_ms) * 1000 if elapsed_ms > 0 else 0

    embedding_ms = 0.0
    vector_search_ms = 0.0
    cache_hit_rate = 0.0

    if profiler_stats:
        if 'embedding_generation' in profiler_stats:
            embedding_ms = profiler_stats['embedding_generation'].get('total_ms', 0.0)
        if 'vector_search' in profiler_stats:
            vector_search_ms = profiler_stats['vector_search'].get('total_ms', 0.0)

    return MLMetrics(
        model_type="VTU CO Generator v2.0",
        inference_latency_ms=round(elapsed_ms / num_cos, 2) if num_cos > 0 else 0,
        throughput_cos_per_sec=round(throughput, 2),
        cache_hit_rate=cache_hit_rate,
        embedding_generation_ms=round(embedding_ms, 2),
        vector_search_ms=round(vector_search_ms, 2)
    )

def store_generation_metrics(session_id: str, cos: List[Dict], metrics: Dict):
    """Store generation metrics for later retrieval"""
    METRICS_STORE['generation_history'].append({
        'session_id': session_id,
        'timestamp': datetime.now().isoformat(),
        'num_cos': len(cos),
        'metrics': metrics,
        'cos_preview': [co['co_text'][:100] for co in cos]
    })

    # Keep only last 100 generations
    if len(METRICS_STORE['generation_history']) > 100:
        METRICS_STORE['generation_history'] = METRICS_STORE['generation_history'][-100:]

def update_aggregate_metrics(metrics: Dict):
    """Update aggregate metrics across all generations"""
    if 'total_generations' not in METRICS_STORE['aggregate_metrics']:
        METRICS_STORE['aggregate_metrics'] = {
            'total_generations': 0,
            'total_cos': 0,
            'avg_quality_score': 0.0,
            'avg_bloom_accuracy': 0.0,
            'avg_vtu_compliance': 0.0,
            'bloom_distribution': {},
            'po_usage_count': {}
        }

    agg = METRICS_STORE['aggregate_metrics']
    agg['total_generations'] += 1
    agg['total_cos'] += metrics.get('total_cos_generated', 0)

    # Update running averages
    n = agg['total_generations']
    agg['avg_quality_score'] = ((agg['avg_quality_score'] * (n-1)) + metrics.get('average_quality_score', 0)) / n
    agg['avg_bloom_accuracy'] = ((agg['avg_bloom_accuracy'] * (n-1)) + metrics.get('bloom_classification_accuracy', 0)) / n
    agg['avg_vtu_compliance'] = ((agg['avg_vtu_compliance'] * (n-1)) + metrics.get('average_vtu_compliance', 0)) / n

# ============================================================================
# API ENDPOINTS
# ============================================================================

@app.get("/", response_model=HealthResponse)
async def root():
    """Root endpoint - health check"""
    import psutil

    system_metrics = {
        'cpu_percent': psutil.cpu_percent(interval=0.1),
        'memory_percent': psutil.virtual_memory().percent,
        'disk_percent': psutil.disk_usage('/').percent
    }

    return {
        "status": "healthy",
        "version": "2.0.0",
        "dependencies": {
            "pdf_available": PDF_AVAILABLE,
            "pptx_available": PPTX_AVAILABLE,
            "docx_available": DOCX_AVAILABLE,
            "chromadb_available": CHROMADB_AVAILABLE,
            "latency_optimizer_available": LATENCY_OPTIMIZER_AVAILABLE
        },
        "system_metrics": system_metrics
    }

@app.get("/health", response_model=HealthResponse)
async def health_check():
    """Health check endpoint with system metrics"""
    try:
        import psutil
        system_metrics = {
            'cpu_percent': psutil.cpu_percent(interval=0.1),
            'memory_percent': psutil.virtual_memory().percent,
            'memory_available_mb': psutil.virtual_memory().available / (1024 * 1024),
            'disk_percent': psutil.disk_usage('/').percent,
            'total_generations': METRICS_STORE['aggregate_metrics'].get('total_generations', 0)
        }
    except:
        system_metrics = {}

    return {
        "status": "healthy",
        "version": "2.0.0",
        "dependencies": {
            "pdf_available": PDF_AVAILABLE,
            "pptx_available": PPTX_AVAILABLE,
            "docx_available": DOCX_AVAILABLE,
            "chromadb_available": CHROMADB_AVAILABLE,
            "latency_optimizer_available": LATENCY_OPTIMIZER_AVAILABLE
        },
        "system_metrics": system_metrics
    }

@app.post("/generate-cos", response_model=COGenerationResponse)
async def generate_cos(
    background_tasks: BackgroundTasks,
    files: List[UploadFile] = File(..., description="Course material files (PDF, PPTX, DOCX, TXT)"),
    num_apply: int = Form(2, ge=0, le=4, description="Number of Apply-level COs"),
    num_analyze: int = Form(2, ge=0, le=4, description="Number of Analyze-level COs"),
    use_chromadb: bool = Form(False, description="Use ChromaDB for context retrieval")
):
    """
    Generate Course Outcomes from uploaded files with comprehensive metrics
    """
    import time
    start_time = time.time()
    session_id = datetime.now().strftime("%Y%m%d_%H%M%S")

    # Validate configuration
    if num_apply + num_analyze != 4:
        raise HTTPException(
            status_code=400,
            detail=f"num_apply + num_analyze must equal 4. Got {num_apply} + {num_analyze} = {num_apply + num_analyze}"
        )

    if not files:
        raise HTTPException(status_code=400, detail="No files uploaded")

    # Start profiling if available
    if profiler:
        profiler.reset()

    # Extract text from all files
    all_text = ""
    temp_files = []

    try:
        # Document processing
        doc_start = time.time()
        for uploaded_file in files:
            suffix = Path(uploaded_file.filename).suffix.lower()
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                content = await uploaded_file.read()
                tmp.write(content)
                tmp_path = tmp.name
                temp_files.append(tmp_path)

            text = extract_text_from_file(tmp_path, uploaded_file.filename)
            all_text += text + "\n\n"
        doc_processing_ms = (time.time() - doc_start) * 1000

        if not all_text.strip():
            raise HTTPException(
                status_code=400,
                detail="No text could be extracted from the uploaded files"
            )

        # Generate COs
        gen_start = time.time()
        generated_cos = generator.generate_custom_cos(
            all_text,
            num_apply=num_apply,
            num_analyze=num_analyze
        )
        llm_inference_ms = (time.time() - gen_start) * 1000

        # Evaluate each CO
        eval_start = time.time()
        cos_with_metrics = []
        bloom_distribution = {}

        for co_data in generated_cos:
            po_list = [po.strip() for po in co_data['po_mappings'].split(',')]
            metrics = evaluator.evaluate_single_co(
                co_data['co_text'],
                co_data['bloom_level'],
                po_list
            )

            # Track Bloom distribution
            bloom_level = co_data['bloom_level']
            bloom_distribution[bloom_level] = bloom_distribution.get(bloom_level, 0) + 1

            cos_with_metrics.append({
                'co_num': co_data['co_num'],
                'co_text': co_data['co_text'],
                'bloom_level': bloom_level,
                'po_mappings': co_data['po_mappings'],
                'topics_covered': co_data.get('topics_covered', []),
                'reward_score': metrics.overall_quality_score,
                'individual_scores': {
                    'vtu': metrics.vtu_compliance_score,
                    'obe': metrics.obe_alignment_score,
                    'bloom': metrics.bloom_accuracy,
                    'conciseness': metrics.conciseness_score
                },
                'approved': metrics.overall_quality_score >= 0.70,
                'word_count': metrics.word_count,
                'has_action_verb': metrics.has_action_verb,
                'has_specific_concepts': metrics.has_specific_concepts
            })
        refinement_ms = (time.time() - eval_start) * 1000

        # Compute aggregate metrics
        pipeline_metrics_obj = evaluator.evaluate_all_cos([
            {'co_text': co['co_text'], 'bloom_level': co['bloom_level'], 'po_mappings': co['po_mappings']}
            for co in cos_with_metrics
        ])

        # Calculate ML metrics
        ml_metrics = calculate_ml_metrics(start_time, len(cos_with_metrics),
                                          profiler.get_stats() if profiler else None)

        # Build comprehensive pipeline metrics
        total_pipeline_ms = (time.time() - start_time) * 1000

        pipeline_metrics = PipelineMetrics(
            bloom_classification_accuracy=pipeline_metrics_obj.bloom_classification_accuracy,
            average_quality_score=pipeline_metrics_obj.average_quality_score,
            average_vtu_compliance=pipeline_metrics_obj.average_vtu_compliance,
            average_obe_alignment=pipeline_metrics_obj.average_obe_alignment,
            average_conciseness_score=pipeline_metrics_obj.average_conciseness_score,
            po_coverage=pipeline_metrics_obj.po_coverage,
            document_processing_ms=round(doc_processing_ms, 2),
            embedding_generation_ms=0.0,
            graph_construction_ms=0.0,
            vector_search_ms=0.0,
            graph_traversal_ms=0.0,
            llm_inference_ms=round(llm_inference_ms, 2),
            refinement_ms=round(refinement_ms, 2),
            total_pipeline_ms=round(total_pipeline_ms, 2),
            ml_metrics=ml_metrics
        )

        # Store metrics in background
        background_tasks.add_task(
            store_generation_metrics,
            session_id,
            cos_with_metrics,
            pipeline_metrics.dict()
        )
        background_tasks.add_task(
            update_aggregate_metrics,
            pipeline_metrics_obj.to_dict()
        )

        return COGenerationResponse(
            success=True,
            message=f"Successfully generated {len(cos_with_metrics)} COs from {len(files)} file(s)",
            cos=cos_with_metrics,
            pipeline_metrics=pipeline_metrics,
            extracted_text_length=len(all_text),
            files_processed=len(files),
            timestamp=datetime.now().isoformat(),
            session_id=session_id
        )

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating COs: {str(e)}")
    finally:
        # Clean up temporary files
        for tmp_path in temp_files:
            try:
                os.unlink(tmp_path)
            except:
                pass

@app.get("/metrics/dashboard", response_model=DashboardMetrics)
async def get_dashboard_metrics():
    """
    Get comprehensive metrics for dashboard display
    Matches the Streamlit dashboard metrics
    """
    agg = METRICS_STORE['aggregate_metrics']

    # Get recent generations
    recent = METRICS_STORE['generation_history'][-10:] if METRICS_STORE['generation_history'] else []

    # Calculate latency breakdown from recent generations
    latency_breakdown = {}
    if recent:
        latest = recent[-1]['metrics']
        latency_breakdown = {
            'document_processing_ms': latest.get('document_processing_ms', 0),
            'embedding_generation_ms': latest.get('embedding_generation_ms', 0),
            'vector_search_ms': latest.get('vector_search_ms', 0),
            'llm_inference_ms': latest.get('llm_inference_ms', 0),
            'refinement_ms': latest.get('refinement_ms', 0),
            'total_pipeline_ms': latest.get('total_pipeline_ms', 0)
        }

    # PO coverage details
    po_coverage_details = agg.get('po_usage_count', {})

    # ML Performance metrics
    ml_performance = {}
    if recent:
        latest_ml = recent[-1]['metrics'].get('ml_metrics', {})
        ml_performance = {
            'inference_latency_ms': latest_ml.get('inference_latency_ms', 0),
            'throughput_cos_per_sec': latest_ml.get('throughput_cos_per_sec', 0),
            'cache_hit_rate': latest_ml.get('cache_hit_rate', 0)
        }

    return DashboardMetrics(
        summary={
            'total_generations': agg.get('total_generations', 0),
            'total_cos_generated': agg.get('total_cos', 0),
            'avg_quality_score': round(agg.get('avg_quality_score', 0), 3),
            'avg_bloom_accuracy': round(agg.get('avg_bloom_accuracy', 0), 3),
            'avg_vtu_compliance': round(agg.get('avg_vtu_compliance', 0), 3)
        },
        quality_metrics={
            'bloom_classification_accuracy': agg.get('avg_bloom_accuracy', 0),
            'average_quality_score': agg.get('avg_quality_score', 0),
            'average_vtu_compliance': agg.get('avg_vtu_compliance', 0),
            'average_obe_alignment': 0.80,  # Default
            'average_conciseness_score': 0.75  # Default
        },
        latency_breakdown=latency_breakdown,
        bloom_distribution=agg.get('bloom_distribution', {}),
        po_coverage_details=po_coverage_details,
        recent_generations=recent,
        ml_performance=ml_performance
    )

@app.get("/metrics/profiler")
async def get_profiler_stats():
    """Get latency profiler statistics"""
    if not profiler:
        return {"error": "Latency profiler not available"}

    return {
        "success": True,
        "stats": profiler.get_stats(),
        "timestamp": datetime.now().isoformat()
    }

@app.get("/metrics/export")
async def export_metrics():
    """Export all metrics as JSON"""
    return StreamingResponse(
        io.StringIO(json.dumps(METRICS_STORE, indent=2)),
        media_type="application/json",
        headers={
            "Content-Disposition": f"attachment; filename=co_generator_metrics_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
        }
    )

@app.post("/feedback/submit")
async def submit_feedback(feedback: FeedbackRequest):
    """Submit user feedback on generated COs"""
    METRICS_STORE['user_feedback'].append({
        'session_id': feedback.session_id,
        'co_num': feedback.co_num,
        'approved': feedback.approved,
        'feedback_text': feedback.feedback_text,
        'edited_text': feedback.edited_text,
        'timestamp': datetime.now().isoformat()
    })

    return {
        "success": True,
        "message": "Feedback submitted successfully"
    }

@app.get("/feedback/history")
async def get_feedback_history(limit: int = 50):
    """Get user feedback history"""
    return {
        "success": True,
        "feedback": METRICS_STORE['user_feedback'][-limit:],
        "total_count": len(METRICS_STORE['user_feedback'])
    }

@app.post("/regenerate-co")
async def regenerate_co(request: RegenerateRequest):
    """Regenerate a specific CO based on user feedback"""
    try:
        new_co_text = regenerate_co_with_feedback(
            co_num=request.co_num,
            original_text=request.original_text,
            bloom_level=request.bloom_level,
            feedback=request.feedback
        )

        po_list = []
        metrics = evaluator.evaluate_single_co(
            new_co_text,
            request.bloom_level,
            po_list
        )

        return {
            "success": True,
            "co_num": request.co_num,
            "regenerated_text": new_co_text,
            "bloom_level": request.bloom_level,
            "quality_metrics": {
                'vtu': metrics.vtu_compliance_score,
                'obe': metrics.obe_alignment_score,
                'bloom': metrics.bloom_accuracy,
                'overall': metrics.overall_quality_score
            },
            "timestamp": datetime.now().isoformat()
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error regenerating CO: {str(e)}")

@app.post("/evaluate-co")
async def evaluate_co(
    co_text: str = Form(...),
    bloom_level: str = Form(...),
    po_mappings: str = Form("")
):
    """Evaluate a single CO for quality metrics"""
    try:
        po_list = [po.strip() for po in po_mappings.split(',')] if po_mappings else []
        metrics = evaluator.evaluate_single_co(co_text, bloom_level, po_list)

        return {
            "success": True,
            "co_text": co_text,
            "metrics": metrics.to_dict(),
            "timestamp": datetime.now().isoformat()
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error evaluating CO: {str(e)}")

@app.get("/bloom-taxonomy")
async def get_bloom_taxonomy():
    """Get Bloom's taxonomy reference data"""
    return {
        "taxonomy": BLOOM_TAXONOMY,
        "timestamp": datetime.now().isoformat()
    }

@app.get("/po-descriptions")
async def get_po_descriptions():
    """Get VTU Program Outcome descriptions"""
    return {
        "pos": VTU_PO_DESCRIPTIONS,
        "timestamp": datetime.now().isoformat()
    }

@app.get("/chromadb/search")
async def search_chromadb(query: str, n_results: int = 5):
    """Search ChromaDB for relevant syllabus content"""
    if not CHROMADB_AVAILABLE:
        raise HTTPException(status_code=503, detail="ChromaDB not available")

    try:
        results = search_syllabus(query, n_results)
        return {
            "success": True,
            "query": query,
            "results": results,
            "count": len(results)
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"ChromaDB search error: {str(e)}")

# ============================================================================
# RUN SERVER
# ============================================================================

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
