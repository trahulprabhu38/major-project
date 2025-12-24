"""
FastAPI Server with Complete 5-Stage Pipeline
==============================================
Uses the full GraphRAG architecture with:
- Document Intelligence
- Knowledge Graph Construction
- Graph-RAG Retrieval
- Multi-Task LLM Generation
- Refinement Layer

With proper CO storage, explainability, and metrics.
"""

from fastapi import FastAPI, File, UploadFile, HTTPException, Form, BackgroundTasks
from fastapi.responses import JSONResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
from typing import List, Optional, Dict, Any
import tempfile
import os
import sys
from pathlib import Path
from datetime import datetime
import json
import io

# Add src to path
sys.path.insert(0, os.path.dirname(__file__))

# Import complete pipeline
from complete_pipeline import CompleteCOPipeline, COStorage

# Import document processing utils
try:
    from PyPDF2 import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

# ============================================================================
# FASTAPI APP INITIALIZATION
# ============================================================================

app = FastAPI(
    title="CO Generator API - Complete Pipeline",
    description="Advanced Multi-Stage AI Pipeline with GraphRAG for VTU-Aligned Course Outcomes",
    version="2.0.0",
    docs_url="/docs",
    redoc_url="/redoc"
)

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize pipeline (singleton)
pipeline = None
co_storage = COStorage()


def get_pipeline():
    """Get or initialize pipeline"""
    global pipeline
    if pipeline is None:
        pipeline = CompleteCOPipeline(
            neo4j_uri=os.getenv("NEO4J_URI", "bolt://localhost:7687"),
            neo4j_user=os.getenv("NEO4J_USER", "neo4j"),
            neo4j_password=os.getenv("NEO4J_PASSWORD", "cogenerator123"),
            chroma_path=os.getenv("CHROMA_PATH", "data/chroma_db")
        )
    return pipeline


# ============================================================================
# PYDANTIC MODELS
# ============================================================================

class COGenerationRequest(BaseModel):
    """Request model for CO generation"""
    num_apply: int = Field(2, ge=0, le=4, description="Number of Apply-level COs")
    num_analyze: int = Field(2, ge=0, le=4, description="Number of Analyze-level COs")
    subject: str = Field("DBMS", description="Subject name")


class CourseOutcome(BaseModel):
    """Model for a single Course Outcome"""
    co_num: int
    co_text: str
    bloom_level: str
    po_mappings: str
    topics_covered: Optional[List[str]] = []
    reward_score: Optional[float] = 0.0
    individual_scores: Optional[Dict[str, float]] = {}
    approved: Optional[bool] = False
    justification: Optional[Dict] = {}


class GraphStats(BaseModel):
    """Knowledge Graph Statistics"""
    nodes: int
    relationships: int
    paths: int
    export_path: Optional[str] = None


class DocumentStats(BaseModel):
    """Document Processing Statistics"""
    files_processed: int
    total_chunks: int


class PipelineMetrics(BaseModel):
    """Comprehensive Pipeline Metrics"""
    bloom_classification_accuracy: float
    average_quality_score: float
    average_vtu_compliance: float
    average_obe_alignment: float
    average_conciseness_score: float
    po_coverage: float


class LatencyMetrics(BaseModel):
    """Latency Metrics"""
    total_ms: float
    avg_per_co_ms: float


class COGenerationResponse(BaseModel):
    """Response model for CO generation"""
    success: bool
    message: str
    session_id: str
    cos: List[CourseOutcome]
    metrics: Dict[str, Any]
    graph_stats: GraphStats
    document_stats: DocumentStats
    latency: LatencyMetrics
    timestamp: str


class HealthResponse(BaseModel):
    """Health check response"""
    status: str
    version: str
    dependencies: Dict[str, bool]
    neo4j_connected: bool


class SessionResponse(BaseModel):
    """Session retrieval response"""
    success: bool
    session: Optional[Dict] = None
    message: Optional[str] = None


# ============================================================================
# TEXT EXTRACTION FUNCTIONS
# ============================================================================

def extract_pdf(file_path: str) -> str:
    """Extract text from PDF"""
    if not PDF_AVAILABLE:
        return ""
    try:
        reader = PdfReader(file_path)
        return "\n".join([page.extract_text() or "" for page in reader.pages])
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading PDF: {e}")


def extract_pptx(file_path: str) -> str:
    """Extract text from PPTX"""
    if not PPTX_AVAILABLE:
        return ""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading PPTX: {e}")


def extract_docx(file_path: str) -> str:
    """Extract text from DOCX"""
    if not DOCX_AVAILABLE:
        return ""
    try:
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading DOCX: {e}")


def extract_text_from_file(file_path: str, filename: str) -> str:
    """Extract text from uploaded file"""
    suffix = Path(filename).suffix.lower()

    try:
        if suffix == '.pdf':
            return extract_pdf(file_path)
        elif suffix in ['.ppt', '.pptx']:
            return extract_pptx(file_path)
        elif suffix in ['.doc', '.docx']:
            return extract_docx(file_path)
        elif suffix == '.txt':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        else:
            return ""
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error extracting text from {filename}: {e}")


# ============================================================================
# API ENDPOINTS
# ============================================================================

@app.get("/", response_model=HealthResponse)
async def root():
    """Root endpoint - health check"""
    pipe = get_pipeline()
    return {
        "status": "healthy",
        "version": "2.0.0",
        "dependencies": {
            "pdf_available": PDF_AVAILABLE,
            "pptx_available": PPTX_AVAILABLE,
            "docx_available": DOCX_AVAILABLE,
        },
        "neo4j_connected": pipe.knowledge_graph.connected if pipe else False
    }


@app.get("/health", response_model=HealthResponse)
async def health_check():
    """Health check endpoint"""
    pipe = get_pipeline()
    return {
        "status": "healthy",
        "version": "2.0.0",
        "dependencies": {
            "pdf_available": PDF_AVAILABLE,
            "pptx_available": PPTX_AVAILABLE,
            "docx_available": DOCX_AVAILABLE,
        },
        "neo4j_connected": pipe.knowledge_graph.connected if pipe else False
    }


@app.post("/generate-cos", response_model=COGenerationResponse)
async def generate_cos(
    background_tasks: BackgroundTasks,
    files: List[UploadFile] = File(..., description="Course material files (PDF, PPTX, DOCX, TXT)"),
    num_apply: int = Form(2, ge=0, le=4, description="Number of Apply-level COs"),
    num_analyze: int = Form(2, ge=0, le=4, description="Number of Analyze-level COs"),
    subject: str = Form("DBMS", description="Subject name")
):
    """
    Generate Course Outcomes using complete 5-stage pipeline

    - Stage 1: Document Intelligence (extract, chunk, embed)
    - Stage 2: Knowledge Graph Construction
    - Stage 3: Graph-RAG Retrieval (vector + graph)
    - Stage 4: Multi-Task LLM Generation
    - Stage 5: Refinement & Reward Scoring
    """

    # Validate configuration
    if num_apply + num_analyze != 4:
        raise HTTPException(
            status_code=400,
            detail=f"num_apply + num_analyze must equal 4. Got {num_apply} + {num_analyze} = {num_apply + num_analyze}"
        )

    if not files:
        raise HTTPException(status_code=400, detail="No files uploaded")

    # Get pipeline
    pipe = get_pipeline()

    # Save uploaded files temporarily
    temp_files = []

    try:
        for uploaded_file in files:
            suffix = Path(uploaded_file.filename).suffix.lower()
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                content = await uploaded_file.read()
                tmp.write(content)
                tmp_path = tmp.name
                temp_files.append(tmp_path)

        # Run complete pipeline
        result = pipe.run_complete_pipeline(
            file_paths=temp_files,
            num_apply=num_apply,
            num_analyze=num_analyze,
            subject=subject
        )

        # Convert to response format
        cos_response = []
        for co in result['cos']:
            cos_response.append(CourseOutcome(
                co_num=int(co['co_text'].split()[0][2:]),  # Extract CO number
                co_text=co['co_text'],
                bloom_level=co['bloom_level'],
                po_mappings=co['po_mappings'],
                topics_covered=co.get('topics_covered', []),
                reward_score=co['scores']['final_score'],
                individual_scores={
                    'conciseness': co['scores']['conciseness'],
                    'vtu_compliance': co['scores']['vtu_compliance'],
                    'obe_alignment': co['scores']['obe_alignment'],
                    'bloom_accuracy': co['scores']['bloom_accuracy'],
                    'specificity': co['scores']['specificity']
                },
                approved=co['approved'],
                justification=co.get('justification', {})
            ))

        return COGenerationResponse(
            success=True,
            message=f"Successfully generated 6 COs from {len(files)} file(s) using complete 5-stage pipeline",
            session_id=result['session_id'],
            cos=cos_response,
            metrics=result['metrics'],
            graph_stats=GraphStats(**result['graph_stats']),
            document_stats=DocumentStats(**result['document_stats']),
            latency=LatencyMetrics(**result['latency']),
            timestamp=datetime.now().isoformat()
        )

    except Exception as e:
        import traceback
        print(f"Error in CO generation: {e}")
        print(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Error generating COs: {str(e)}")
    finally:
        # Clean up temporary files
        for tmp_path in temp_files:
            try:
                os.unlink(tmp_path)
            except:
                pass


@app.get("/sessions")
async def list_sessions():
    """List all saved CO generation sessions"""
    sessions = co_storage.get_all_sessions()
    return {
        "success": True,
        "total_sessions": len(sessions),
        "sessions": sessions
    }


@app.get("/sessions/{session_id}", response_model=SessionResponse)
async def get_session(session_id: str):
    """Get a specific session by ID"""
    session = co_storage.get_session(session_id)
    if session:
        return SessionResponse(success=True, session=session)
    else:
        return SessionResponse(
            success=False,
            message=f"Session {session_id} not found"
        )


@app.get("/sessions/{session_id}/export")
async def export_session(session_id: str):
    """Export a session as JSON file"""
    session = co_storage.get_session(session_id)
    if not session:
        raise HTTPException(status_code=404, detail=f"Session {session_id} not found")

    return StreamingResponse(
        io.StringIO(json.dumps(session, indent=2)),
        media_type="application/json",
        headers={
            "Content-Disposition": f"attachment; filename=session_{session_id}.json"
        }
    )


@app.get("/statistics")
async def get_statistics():
    """Get overall system statistics"""
    stats = co_storage.get_statistics()
    pipe = get_pipeline()

    return {
        "success": True,
        "storage_stats": stats,
        "pipeline_status": {
            "neo4j_connected": pipe.knowledge_graph.connected,
            "chromadb_ready": pipe.graph_rag.vector_db_ready
        }
    }


@app.get("/knowledge-graph")
async def get_knowledge_graph():
    """Get knowledge graph data"""
    pipe = get_pipeline()

    return {
        "success": True,
        "graph_data": pipe.knowledge_graph.graph_data,
        "statistics": {
            "nodes": len(pipe.knowledge_graph.graph_data['nodes']),
            "relationships": len(pipe.knowledge_graph.graph_data['relationships']),
            "paths": len(pipe.knowledge_graph.graph_data['paths'])
        }
    }


@app.post("/reset-graph")
async def reset_knowledge_graph():
    """Reset the knowledge graph (useful for testing)"""
    global pipeline
    pipeline = None  # Force re-initialization
    return {
        "success": True,
        "message": "Knowledge graph will be rebuilt on next request"
    }


@app.get("/sessions/{session_id}/co-po-mapping")
async def get_co_po_mapping(session_id: str):
    """
    Get CO-PO mapping matrix for a specific session
    Returns data formatted for frontend visualization
    """
    session = co_storage.get_session(session_id)
    if not session:
        raise HTTPException(status_code=404, detail=f"Session {session_id} not found")

    # Extract CO-PO mapping data
    cos = session.get('cos', [])

    # Build mapping matrix
    po_list = [f"PO{i}" for i in range(1, 13)]  # PO1 to PO12

    mapping_matrix = []
    co_details = []

    for co in cos:
        co_num = int(co['co_text'].split()[0][2:])  # Extract CO number
        co_text = co['co_text']
        bloom_level = co['bloom_level']
        po_mappings = co.get('po_mappings', '')

        # Parse PO mappings
        mapped_pos = [po.strip() for po in po_mappings.split(',') if po.strip()]

        # Create row for matrix (1 if mapped, 0 if not)
        row = {
            'co_num': co_num,
            'mappings': {po: (1 if po in mapped_pos else 0) for po in po_list}
        }
        mapping_matrix.append(row)

        # Store CO details
        co_details.append({
            'co_num': co_num,
            'co_text': co_text,
            'bloom_level': bloom_level,
            'po_mappings': mapped_pos,
            'score': co.get('scores', {}).get('final_score', 0)
        })

    # Calculate PO coverage statistics
    po_coverage = {po: 0 for po in po_list}
    for co in cos:
        mapped_pos = [po.strip() for po in co.get('po_mappings', '').split(',') if po.strip()]
        for po in mapped_pos:
            if po in po_coverage:
                po_coverage[po] += 1

    return {
        "success": True,
        "session_id": session_id,
        "co_po_matrix": {
            "headers": {
                "cos": [f"CO{i}" for i in range(1, 7)],
                "pos": po_list
            },
            "matrix": mapping_matrix,
            "co_details": co_details
        },
        "statistics": {
            "total_cos": len(cos),
            "total_pos": len(po_list),
            "po_coverage": po_coverage,
            "coverage_percentage": {
                po: round((count / len(cos)) * 100, 1) if len(cos) > 0 else 0
                for po, count in po_coverage.items()
            }
        }
    }


@app.get("/co-po-mapping/latest")
async def get_latest_co_po_mapping():
    """
    Get CO-PO mapping for the most recent session
    Convenient endpoint for frontend to fetch latest results
    """
    sessions = co_storage.get_all_sessions()
    if not sessions:
        raise HTTPException(status_code=404, detail="No sessions found")

    latest_session = sessions[-1]
    session_id = latest_session['session_id']

    # Reuse the mapping logic
    return await get_co_po_mapping(session_id)


@app.get("/co-po-mapping/visualize/{session_id}")
async def visualize_co_po_mapping(session_id: str):
    """
    Get CO-PO mapping in format optimized for chart libraries
    Returns data ready for heatmaps, charts, and tables
    """
    session = co_storage.get_session(session_id)
    if not session:
        raise HTTPException(status_code=404, detail=f"Session {session_id} not found")

    cos = session.get('cos', [])
    po_list = [f"PO{i}" for i in range(1, 13)]

    # Format 1: Heatmap data (for libraries like Chart.js, Plotly)
    heatmap_data = {
        'x': po_list,  # X-axis: POs
        'y': [f"CO{i}" for i in range(1, 7)],  # Y-axis: COs
        'z': []  # 2D array: values
    }

    for co in cos:
        co_num = int(co['co_text'].split()[0][2:])
        po_mappings = co.get('po_mappings', '')
        mapped_pos = [po.strip() for po in po_mappings.split(',') if po.strip()]

        row = [1 if po in mapped_pos else 0 for po in po_list]
        heatmap_data['z'].append(row)

    # Format 2: Table data (for DataTables, AG Grid)
    table_data = []
    for co in cos:
        co_num = int(co['co_text'].split()[0][2:])
        co_text = co['co_text']
        bloom_level = co['bloom_level']
        po_mappings = co.get('po_mappings', '')
        mapped_pos = [po.strip() for po in po_mappings.split(',') if po.strip()]

        row = {
            'co': f"CO{co_num}",
            'description': co_text,
            'bloom_level': bloom_level,
            'score': round(co.get('scores', {}).get('final_score', 0), 2)
        }

        # Add PO columns
        for po in po_list:
            row[po] = '✓' if po in mapped_pos else ''

        table_data.append(row)

    # Format 3: Graph/Network data (for D3.js, Cytoscape)
    graph_data = {
        'nodes': [],
        'edges': []
    }

    # Add CO nodes
    for co in cos:
        co_num = int(co['co_text'].split()[0][2:])
        graph_data['nodes'].append({
            'id': f"CO{co_num}",
            'label': f"CO{co_num}",
            'type': 'co',
            'bloom_level': co['bloom_level'],
            'score': co.get('scores', {}).get('final_score', 0)
        })

    # Add PO nodes
    for po in po_list:
        graph_data['nodes'].append({
            'id': po,
            'label': po,
            'type': 'po'
        })

    # Add edges (CO -> PO mappings)
    for co in cos:
        co_num = int(co['co_text'].split()[0][2:])
        po_mappings = co.get('po_mappings', '')
        mapped_pos = [po.strip() for po in po_mappings.split(',') if po.strip()]

        for po in mapped_pos:
            graph_data['edges'].append({
                'source': f"CO{co_num}",
                'target': po,
                'weight': 1
            })

    # Format 4: Summary statistics
    summary = {
        'total_mappings': sum(
            len([po.strip() for po in co.get('po_mappings', '').split(',') if po.strip()])
            for co in cos
        ),
        'avg_mappings_per_co': round(
            sum(len([po.strip() for po in co.get('po_mappings', '').split(',') if po.strip()]) for co in cos) / len(cos), 1
        ) if len(cos) > 0 else 0,
        'bloom_distribution': {},
        'po_utilization': {}
    }

    # Bloom level distribution
    for co in cos:
        bloom = co['bloom_level']
        summary['bloom_distribution'][bloom] = summary['bloom_distribution'].get(bloom, 0) + 1

    # PO utilization
    for po in po_list:
        count = sum(
            1 for co in cos
            if po in [p.strip() for p in co.get('po_mappings', '').split(',')]
        )
        summary['po_utilization'][po] = {
            'count': count,
            'percentage': round((count / len(cos)) * 100, 1) if len(cos) > 0 else 0
        }

    return {
        "success": True,
        "session_id": session_id,
        "formats": {
            "heatmap": heatmap_data,
            "table": table_data,
            "graph": graph_data
        },
        "summary": summary,
        "metadata": {
            "timestamp": session.get('timestamp'),
            "source_files": session.get('metadata', {}).get('source_files', [])
        }
    }


# ============================================================================
# RUN SERVER
# ============================================================================

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
