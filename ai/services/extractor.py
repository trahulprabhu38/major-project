import fitz  # PyMuPDF
import docx
from pptx import Presentation

def extract_text(path: str) -> str:
    text = ""
    if path.endswith(".pdf"):
        pdf = fitz.open(path)
        text = " ".join(page.get_text() for page in pdf)
    elif path.endswith(".docx"):
        doc = docx.Document(path)
        text = " ".join(p.text for p in doc.paragraphs)
    elif path.endswith(".pptx"):
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        text = " ".join(texts)
    return text.strip()
